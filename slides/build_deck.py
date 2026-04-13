"""Build soma presentation deck — light theme, content grounded in docs."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
BG       = RGBColor(0xFF, 0xFF, 0xFF)   # white
SURFACE  = RGBColor(0xF4, 0xF5, 0xF7)  # light gray card
BORDER   = RGBColor(0xD8, 0xDA, 0xDE)  # subtle rule
ACCENT   = RGBColor(0x00, 0x96, 0x6E)  # teal
ACCENT2  = RGBColor(0x2F, 0x6F, 0xC8)  # blue
TEXT     = RGBColor(0x1A, 0x1A, 0x1A)  # near-black
MUTED    = RGBColor(0x66, 0x6B, 0x74)  # gray
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Primitives ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def txt(slide, text, x, y, w, h, size=14, bold=False, italic=False,
        color=TEXT, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def img_placeholder(slide, x, y, w, h, label="[ image placeholder ]"):
    rect(slide, x, y, w, h, SURFACE, BORDER)
    txt(slide, label, x, y + h // 2 - Inches(0.18), w, Inches(0.36),
        size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


def title_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(0.055), ACCENT)
    txt(slide, title,
        Inches(0.6), Inches(0.18), Inches(12.1), Inches(0.72),
        size=30, bold=True, color=TEXT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.6), Inches(0.88), Inches(11), Inches(0.4),
            size=15, color=ACCENT)


def slide_num(slide, n, total=10):
    txt(slide, f"{n} / {total}",
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.38), Inches(1.0), Inches(0.32),
        size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def bullet_card(slide, items, x, y, w, h, item_size=14):
    rect(slide, x, y, w, h, SURFACE)
    txb = slide.shapes.add_textbox(
        x + Inches(0.22), y + Inches(0.18),
        w - Inches(0.44), h - Inches(0.36))
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        run = p.add_run()
        run.text = "·  " + item
        run.font.size  = Pt(item_size)
        run.font.color.rgb = TEXT


def label_card(slide, title, subtitle, x, y, w, h,
               title_color=TEXT, title_size=14):
    rect(slide, x, y, w, h, SURFACE)
    txt(slide, title, x + Inches(0.15), y + Inches(0.08),
        w - Inches(0.3), Inches(0.42),
        size=title_size, bold=True, color=title_color)
    txt(slide, subtitle, x + Inches(0.15), y + Inches(0.46),
        w - Inches(0.3), h - Inches(0.54),
        size=11, color=MUTED)


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    s = blank_slide(prs)
    fill_bg(s)
    rect(s, 0, 0, Inches(0.1), SLIDE_H, ACCENT)
    txt(s, "soma",
        Inches(0.65), Inches(1.6), Inches(7.5), Inches(1.8),
        size=80, bold=True, color=TEXT)
    txt(s, "slides  ·  labels  →  trained model  ·  evaluation report",
        Inches(0.65), Inches(3.45), Inches(8.5), Inches(0.55),
        size=18, color=ACCENT)
    txt(s, "A modular, reproducible framework for computational pathology research",
        Inches(0.65), Inches(4.1), Inches(8.5), Inches(0.7),
        size=15, italic=True, color=MUTED)
    txt(s, "Clément Grisi",
        Inches(0.65), SLIDE_H - Inches(1.0), Inches(5), Inches(0.4),
        size=13, color=MUTED)
    img_placeholder(s, Inches(9.4), Inches(1.4), Inches(3.5), Inches(4.8),
                    "[ pipeline overview ]")


def slide_origin(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "01 — The Starting Point",
              "Automating FM evaluation in computational pathology")

    bullet_card(s, [
        "Needed an automated evaluation pipeline for pathology foundation models",
        "Preprocessing delegated to hs2p  (tissue masking, tiling, coordinates)",
        "Encoding delegated to slide2vec  (WSI reader, batch inference, embeddings)",
        "Evaluation logic added on top  →  evaluation reports + per-case predictions",
    ], Inches(0.6), Inches(1.48), Inches(7.6), Inches(3.3))

    # dependency box
    for i, (name, role) in enumerate([
        ("hs2p",       "tiling & tissue masking"),
        ("slide2vec",  "encoding & WSI reader stack"),
        ("soma",       "training + evaluation"),
    ]):
        col = ACCENT if name == "soma" else ACCENT2
        bx = Inches(8.6)
        by = Inches(1.8) + i * Inches(1.05)
        rect(s, bx, by, Inches(4.3), Inches(0.8), SURFACE)
        rect(s, bx, by, Inches(0.06), Inches(0.8), col)
        txt(s, name, bx + Inches(0.2), by + Inches(0.06),
            Inches(2.0), Inches(0.38), size=15, bold=True, color=col)
        txt(s, role, bx + Inches(0.2), by + Inches(0.42),
            Inches(3.9), Inches(0.3), size=12, color=MUTED)

    slide_num(s, 1)


def slide_scope(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "02 — Scope Creep (The Good Kind)",
              "From an evaluation helper to the full compath pipeline")

    bullet_card(s, [
        "Added linear probing for slide-level FMs",
        "Added MIL aggregators for tile-level FMs",
        "Realisation: this is the full compath pipeline end-to-end",
    ], Inches(0.6), Inches(1.48), Inches(12.1), Inches(1.8))

    # pipeline row
    stages = [
        ("Slides",        "input"),
        ("Preprocessing", "hs2p"),
        ("Encoding",      "slide2vec"),
        ("Aggregation",   "soma MIL"),
        ("Prediction",    "soma heads"),
        ("Report",        "soma evaluation"),
    ]
    colors = [BORDER, ACCENT2, ACCENT2, ACCENT, ACCENT, ACCENT]
    bw, bh = Inches(1.78), Inches(0.88)
    gap    = Inches(0.14)
    sx     = Inches(0.6)
    sy     = Inches(3.6)
    for i, ((name, sub), col) in enumerate(zip(stages, colors)):
        x = sx + i * (bw + gap)
        rect(s, x, sy, bw, bh, SURFACE)
        rect(s, x, sy, bw, Inches(0.055), col)
        txt(s, name, x, sy + Inches(0.1), bw, Inches(0.42),
            size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        txt(s, sub, x, sy + Inches(0.5), bw, Inches(0.32),
            size=11, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            txt(s, "→", x + bw, sy + Inches(0.22), gap + Inches(0.05), Inches(0.44),
                size=14, color=MUTED, align=PP_ALIGN.CENTER)

    slide_num(s, 2)


def slide_vision(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "03 — The Auto-Research Vision",
              "soma as the execution engine for agent-driven experiment search")

    bullet_card(s, [
        "Inspired by Karpathy's auto-research tooling",
        "Target problems: ISUP grading from biopsies, BCR risk from prostatectomies",
        "Agents need a clean, programmable API — not a training script collection",
        "soma exposes:  Pipeline(config).run()  ·  train(store, ...)  ·  FeatureExtractor.extract()",
    ], Inches(0.6), Inches(1.48), Inches(7.6), Inches(3.3))

    # two-layer stack
    for i, (label, col, bg) in enumerate([
        ("🤖  Auto-Research Agent\n(sweep · propose · compare)", ACCENT,  RGBColor(0xE6, 0xF5, 0xF0)),
        ("soma  —  execution layer\n(run · train · evaluate · log)",  ACCENT2, RGBColor(0xE8, 0xF0, 0xFB)),
    ]):
        bx = Inches(8.6)
        by = Inches(1.75) + i * Inches(1.6)
        rect(s, bx, by, Inches(4.3), Inches(1.2), bg)
        rect(s, bx, by, Inches(0.07), Inches(1.2), col)
        txt(s, label, bx + Inches(0.22), by + Inches(0.22),
            Inches(3.9), Inches(0.9), size=14, bold=True, color=col)
        if i == 0:
            txt(s, "↓  calls", Inches(8.6), Inches(2.98), Inches(4.3), Inches(0.38),
                size=13, color=MUTED, align=PP_ALIGN.CENTER)

    slide_num(s, 3)


def slide_landscape(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "04 — Existing Landscape",
              "Why build yet another framework?")

    for i, (name, pts, col) in enumerate([
        ("Slideflow", [
            "Covers the full pipeline end-to-end",
            "No longer maintained",
            "Built on fastai — not native PyTorch",
        ], ACCENT2),
        ("LazySlide", [
            "Modern, well-structured, ambitious scope",
            "Does not train task-specific models on FM outputs",
            "Gap remains for downstream MIL training",
        ], ACCENT2),
    ]):
        x = Inches(0.6) + i * Inches(6.3)
        rect(s, x, Inches(1.5), Inches(5.9), Inches(3.0), SURFACE)
        rect(s, x, Inches(1.5), Inches(5.9), Inches(0.055), col)
        txt(s, name, x + Inches(0.2), Inches(1.6),
            Inches(5.5), Inches(0.5), size=20, bold=True, color=TEXT)
        for j, pt in enumerate(pts):
            txt(s, "✗  " + pt,
                x + Inches(0.2), Inches(2.2) + j * Inches(0.7),
                Inches(5.5), Inches(0.6), size=13, color=MUTED)

    rect(s, Inches(0.6), Inches(4.75), Inches(12.1), Inches(0.95), RGBColor(0xE6, 0xF5, 0xF0))
    rect(s, Inches(0.6), Inches(4.75), Inches(0.07), Inches(0.95), ACCENT)
    txt(s,
        "No well-maintained, PyTorch-native framework trains task-specific models "
        "on top of FM outputs  →  this is soma's space",
        Inches(0.85), Inches(4.82), Inches(11.7), Inches(0.8),
        size=15, bold=True, color=ACCENT)

    slide_num(s, 4)


def slide_intro(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "05 — Introducing soma",
              "The deep learning body for computational pathology")

    cards = [
        ("Reproducible",   "Deterministic splits · versioned configs · content-addressed feature cache"),
        ("Modular",        "Swap encoder, aggregator, or task head without touching unrelated code"),
        ("Agent-friendly", "Pipeline(config).run()  ·  train(store, ...)  ·  FeatureExtractor.extract()"),
        ("Fast",           "cucim batch region reads · near 100 % GPU utilisation during encoding"),
    ]
    bw, bh = Inches(5.8), Inches(1.1)
    for i, (title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.5) + row * Inches(1.25)
        rect(s, x, y, bw, bh, SURFACE)
        rect(s, x, y, Inches(0.06), bh, ACCENT)
        txt(s, title, x + Inches(0.2), y + Inches(0.1),
            Inches(2.5), Inches(0.44), size=15, bold=True, color=ACCENT)
        txt(s, body, x + Inches(0.2), y + Inches(0.54),
            bw - Inches(0.35), Inches(0.5), size=12, color=MUTED)

    img_placeholder(s, Inches(1.5), Inches(4.2), Inches(10.3), Inches(2.0),
                    "[ structured YAML config snapshot ]")
    slide_num(s, 5)


def slide_architecture(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "06 — Pipeline Architecture",
              "Composable components · explicit artifact boundaries")

    components = [
        ("Dataset",       "dataset.csv  ·  splits.csv\npatient-level splits"),
        ("Preprocessor",  "tissue mask  ·  tiling\n→ hs2p / slide2vec"),
        ("Featurizer",    "18 FM encoders\n→ slide2vec runtime"),
        ("Aggregator",    "9 MIL heads\nABMIL · CLAM · TransMIL …"),
        ("Predictor",     "classification · ordinal\nregression"),
        ("Evaluator",     "per-fold metrics\npredictions.csv"),
    ]
    colors = [BORDER, ACCENT2, ACCENT2, ACCENT, ACCENT, ACCENT]
    bw, bh = Inches(1.85), Inches(1.45)
    gap    = Inches(0.14)
    sx     = Inches(0.5)
    sy     = Inches(1.6)
    for i, ((name, sub), col) in enumerate(zip(components, colors)):
        x = sx + i * (bw + gap)
        rect(s, x, sy, bw, bh, SURFACE)
        rect(s, x, sy, bw, Inches(0.055), col)
        txt(s, name, x, sy + Inches(0.1), bw, Inches(0.44),
            size=13, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
        txt(s, sub, x, sy + Inches(0.54), bw, Inches(0.85),
            size=10, color=MUTED, align=PP_ALIGN.CENTER)
        if i < len(components) - 1:
            txt(s, "→", x + bw, sy + Inches(0.5), gap + Inches(0.1), Inches(0.44),
                size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # artifact row
    artifacts = [
        "dataset.csv\nsplits.csv",
        "process_list.csv\ncoordinates/",
        "tile_embeddings/\nhierarchical_embeddings/",
        "bag repr\n(B, D)",
        "logits\nprobabilities",
        "metrics.json\npredictions.csv",
    ]
    ay = sy + bh + Inches(0.35)
    for i, art in enumerate(artifacts):
        x = sx + i * (bw + gap)
        txt(s, art, x, ay, bw, Inches(0.65),
            size=9, color=MUTED, align=PP_ALIGN.CENTER, italic=True)

    txt(s, "artifacts", Inches(0.5), ay, Inches(1.0), Inches(0.32),
        size=9, bold=True, color=BORDER)

    img_placeholder(s, Inches(1.5), Inches(5.25), Inches(10.3), Inches(1.75),
                    "[ output tree:  experiments/<slug>/runs/<run_id>/fold_N/ ]")
    slide_num(s, 6)


def slide_speed(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "07 — Engineering Highlight: Encoding Speed",
              "Near 100 % GPU utilisation — cucim batch region reads")

    bullet_card(s, [
        "Encoding is the dominant runtime bottleneck in most compath workflows",
        "Naive approach: read one tile at a time  →  GPU starved by WSI I/O",
        "soma uses cucim batch region reads: large contiguous reads per slide",
        "Adaptive batching and multi-GPU torchrun orchestration via slide2vec",
        "Result: GPU utilisation ≈ 100 % throughout feature extraction",
    ], Inches(0.6), Inches(1.5), Inches(7.6), Inches(4.0))

    img_placeholder(s, Inches(8.5), Inches(1.5), Inches(4.4), Inches(4.0),
                    "[ GPU util: naive vs soma ]")
    slide_num(s, 7)


def slide_training(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "08 — Downstream Training",
              "Task heads + MIL aggregators · multi-fold CV · per-case predictions")

    # Task heads
    txt(s, "Task Heads", Inches(0.6), Inches(1.5), Inches(5.9), Inches(0.4),
        size=16, bold=True, color=ACCENT)
    heads = [
        ("classification",
         "binary & multi-class · cross-entropy · AUC · F1",
         ACCENT),
        ("ordinal_classification",
         "ordered grades · MSE loss · QWK metric · raw_score saved",
         ACCENT),
        ("regression",
         "continuous targets · MSE · MAE · R²",
         ACCENT),
    ]
    for i, (name, sub, col) in enumerate(heads):
        x, y = Inches(0.6), Inches(1.95) + i * Inches(1.0)
        rect(s, x, y, Inches(5.9), Inches(0.82), SURFACE)
        rect(s, x, y, Inches(0.06), Inches(0.82), col)
        txt(s, name, x + Inches(0.18), y + Inches(0.06),
            Inches(5.5), Inches(0.38), size=13, bold=True, color=TEXT)
        txt(s, sub, x + Inches(0.18), y + Inches(0.44),
            Inches(5.5), Inches(0.32), size=11, color=MUTED)

    # MIL aggregators
    txt(s, "MIL Aggregators", Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.4),
        size=16, bold=True, color=ACCENT2)
    mils = [
        ("abmil",    "Attention-Based MIL  (Ilse 2018)"),
        ("clam_sb",  "Single-branch CLAM, task-aware aux loss  (Lu 2021)"),
        ("clam_mb",  "Multi-branch CLAM  (Lu 2021)"),
        ("dsmil",    "Dual-Stream MIL  (Li 2021)"),
        ("transmil", "Transformer MIL, Nyström attn  (Shao 2021)"),
        ("dtfdmil",  "Double-Tier Feature Distillation  (Zhang 2022)"),
        ("hipt",     "Hierarchical Pyramid Transformer  (Chen 2022)"),
    ]
    for i, (name, sub) in enumerate(mils):
        x, y = Inches(7.0), Inches(1.95) + i * Inches(0.66)
        rect(s, x, y, Inches(5.9), Inches(0.54), SURFACE)
        rect(s, x, y, Inches(0.06), Inches(0.54), ACCENT2)
        txt(s, name, x + Inches(0.18), y + Inches(0.04),
            Inches(1.6), Inches(0.34), size=12, bold=True, color=TEXT)
        txt(s, sub, x + Inches(1.85), y + Inches(0.08),
            Inches(3.85), Inches(0.36), size=11, color=MUTED)

    slide_num(s, 8)


def slide_api(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "09 — Two-Layer API",
              "Composable building blocks  or  single-call pipeline")

    # Layer 1
    rect(s, Inches(0.6), Inches(1.48), Inches(5.9), Inches(4.5), SURFACE)
    rect(s, Inches(0.6), Inches(1.48), Inches(5.9), Inches(0.055), ACCENT2)
    txt(s, "Layer 1 — Step APIs",
        Inches(0.8), Inches(1.55), Inches(5.5), Inches(0.42),
        size=15, bold=True, color=ACCENT2)
    txt(s,
        "Extract once, sweep many models:\n\n"
        "  store = extractor.extract(...)\n\n"
        "  for agg in [\"abmil\", \"clam_sb\",\n"
        "               \"transmil\", ...]:\n"
        "      train(store, dataset, splits,\n"
        "            aggregator=AggregatorConfig(\n"
        "              name=agg), ...)",
        Inches(0.8), Inches(2.05), Inches(5.5), Inches(3.7),
        size=12, color=TEXT)

    # Layer 2
    rect(s, Inches(7.0), Inches(1.48), Inches(5.9), Inches(4.5), SURFACE)
    rect(s, Inches(7.0), Inches(1.48), Inches(5.9), Inches(0.055), ACCENT)
    txt(s, "Layer 2 — Pipeline Orchestrator",
        Inches(7.2), Inches(1.55), Inches(5.5), Inches(0.42),
        size=15, bold=True, color=ACCENT)
    txt(s,
        "Single call end-to-end:\n\n"
        "  result = Pipeline(\n"
        "    PipelineConfig(\n"
        "      dataset_csv=...,\n"
        "      splits_csv=...,\n"
        "      output_root=...,\n"
        "      encoder=EncoderConfig(\n"
        "        name=\"uni2\"),\n"
        "      aggregator=AggregatorConfig(\n"
        "        name=\"abmil\"),\n"
        "      task=TaskConfig(\n"
        "        name=\"classification\"),\n"
        "  )).run()",
        Inches(7.2), Inches(2.05), Inches(5.5), Inches(3.7),
        size=12, color=TEXT)

    slide_num(s, 9)


def slide_roadmap(prs):
    s = blank_slide(prs)
    fill_bg(s)
    title_bar(s, "10 — Roadmap", "What comes next")

    horizons = [
        ("Short-term",  ACCENT,                       [
            "Survival modelling  (Cox, discrete-time)",
            "Attention heatmaps for MIL models",
        ]),
        ("Mid-term",    ACCENT2,                      [
            "Dense task support  (segmentation, detection)",
        ]),
        ("Long-term",   RGBColor(0x7C, 0x4D, 0xBF),  [
            "Auto-derive cell detection / segmentation outputs",
            "Feed as extra signal to AI agent reasoning over dataset + task",
            "Auto-research engine: sweep preprocessing × FM × aggregator",
        ]),
    ]
    y = Inches(1.6)
    for label, col, items in horizons:
        bar_h = Inches(len(items) * 0.68 + 0.1)
        # thin vertical accent bar only — no card background
        rect(s, Inches(0.6), y, Inches(0.1), bar_h, col)
        txt(s, label,
            Inches(0.9), y, Inches(2.2), Inches(0.42),
            size=15, bold=True, color=col)
        for j, item in enumerate(items):
            txt(s, "•  " + item,
                Inches(0.9), y + Inches(0.42) + j * Inches(0.68),
                Inches(11.8), Inches(0.58),
                size=15, color=TEXT)
        y += bar_h + Inches(0.55)

    slide_num(s, 10)


def slide_closing(prs):
    s = blank_slide(prs)
    fill_bg(s)
    rect(s, 0, 0, SLIDE_W, Inches(0.055), ACCENT)
    txt(s, "soma", Inches(0.6), Inches(2.0), Inches(12.1), Inches(1.6),
        size=72, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    txt(s, "reproducible  ·  modular  ·  agent-ready",
        Inches(0.6), Inches(3.65), Inches(12.1), Inches(0.55),
        size=20, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(s, "github.com/cgrisi/soma",
        Inches(0.6), SLIDE_H - Inches(1.0), Inches(12.1), Inches(0.4),
        size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()
    slide_cover(prs)
    slide_origin(prs)
    slide_scope(prs)
    slide_vision(prs)
    slide_landscape(prs)
    slide_intro(prs)
    slide_architecture(prs)
    slide_speed(prs)
    slide_training(prs)
    slide_api(prs)
    slide_roadmap(prs)
    slide_closing(prs)
    out = "slides/soma_presentation.pptx"
    prs.save(out)
    print(f"Saved → {out}")


if __name__ == "__main__":
    main()
